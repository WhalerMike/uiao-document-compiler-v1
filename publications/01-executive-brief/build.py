#!/usr/bin/env python3
# build.py — Automated document pipeline for UIAO Executive Brief
#
# Generates images from prompts, inserts them into the source .docx,
# then produces .docx, .pptx, .pdf, and .epub outputs.
#
# Usage:
#   python build.py                  # Full pipeline: generate images + all outputs
#   python build.py --skip-images    # Skip image generation, use existing images
#   python build.py --formats docx pdf   # Only produce specific formats
#   python build.py --dry-run        # Show what would happen without doing it
#
# Prerequisites:
#   pip install google-genai Pillow python-pptx ebooklib lxml
#   LibreOffice installed (for PDF conversion)

import argparse
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import time
import zipfile
from copy import deepcopy
from pathlib import Path
from lxml import etree

# ──────────────────────────────────────────────────────────────
# CONFIGURATION
# ──────────────────────────────────────────────────────────────
BASE_DIR = Path(__file__).resolve().parent
PROMPTS_FILE = BASE_DIR.parent.parent / "prompts.json"
IMAGES_DIR = BASE_DIR.parent.parent / "images"
SOURCE_DOCX = BASE_DIR / "UIAO-Executive-Brief-SOURCE.docx"
API_KEY = "AIzaSyCQ-BRx8IuLV23ybUjluLY1WfT6mxS9-jQ"
MODEL = "gemini-2.5-flash-image"

# Output filenames
OUTPUT_DOCX = BASE_DIR / "UIAO-Executive-Brief.docx"
OUTPUT_PPTX = BASE_DIR / "UIAO-Executive-Brief.pptx"
OUTPUT_PDF = BASE_DIR / "UIAO-Executive-Brief.pdf"
OUTPUT_EPUB = BASE_DIR / "UIAO-Executive-Brief.epub"

# Image dimensions in EMUs (4 inches wide, square for 1024x1024 source)
# 914400 EMU = 1 inch; page content width ~6.1 inches
IMG_EMU_W = 3657600
IMG_EMU_H = 3657600

# OOXML namespaces
NSMAP = {
    "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "wp": "http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing",
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "pic": "http://schemas.openxmlformats.org/drawingml/2006/picture",
    "ct": "http://schemas.openxmlformats.org/package/2006/content-types",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
# ──────────────────────────────────────────────────────────────


def generate_images(prompts_file, images_dir, api_key, model, delay=2.0):
    """Generate images from prompts using the Gemini API."""
    from google import genai

    with open(prompts_file, "r", encoding="utf-8") as f:
        prompts = json.load(f)

    images_dir.mkdir(parents=True, exist_ok=True)
    client = genai.Client(api_key=api_key)

    print(f"\n=== Generating {len(prompts)} image(s) via {model} ===\n")
    succeeded, failed = 0, 0

    for i, item in enumerate(prompts, 1):
        filename = item["filename"]
        output_path = images_dir / filename

        # Skip if image already exists and is non-empty
        if output_path.exists() and output_path.stat().st_size > 0:
            print(f"  [{i}/{len(prompts)}] Skipping {filename} (already exists)")
            succeeded += 1
            continue

        print(f"  [{i}/{len(prompts)}] Generating: {item['id']} -> {filename}")
        try:
            response = client.models.generate_content(model=model, contents=[item["prompt"]])
            for part in response.parts:
                if part.inline_data is not None:
                    part.as_image().save(str(output_path))
                    size_kb = output_path.stat().st_size / 1024
                    print(f"    Saved ({size_kb:.0f} KB)")
                    succeeded += 1
                    break
            else:
                print(f"    Warning: No image data returned")
                failed += 1
        except Exception as e:
            print(f"    Error: {e}")
            failed += 1

        if i < len(prompts):
            time.sleep(delay)

    print(f"\n  Images: {succeeded} succeeded, {failed} failed\n")
    return failed == 0


def get_image_map(prompts_file, images_dir):
    """Build a mapping from placeholder tag to image file path."""
    with open(prompts_file, "r", encoding="utf-8") as f:
        prompts = json.load(f)

    image_map = {}
    for item in prompts:
        # Extract the number from the id (e.g., "image_05" -> "05")
        num = item["id"].split("_")[1]
        tag = f"IMAGE_{num}"
        img_path = images_dir / item["filename"]
        if img_path.exists():
            image_map[tag] = img_path
        else:
            print(f"  Warning: Image not found for [{tag}]: {img_path}")
    return image_map


def build_drawing_xml(rId, doc_pr_id, name, desc, cx, cy):
    """Build the OOXML drawing element for an inline image."""
    return f'''<w:p xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing" xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:pic="http://schemas.openxmlformats.org/drawingml/2006/picture">
  <w:pPr><w:spacing w:before="240" w:after="240"/><w:jc w:val="center"/></w:pPr>
  <w:r>
    <w:drawing>
      <wp:inline distT="0" distB="0" distL="0" distR="0">
        <wp:extent cx="{cx}" cy="{cy}"/>
        <wp:docPr id="{doc_pr_id}" name="{name}" descr="{desc}"/>
        <a:graphic>
          <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/picture">
            <pic:pic>
              <pic:nvPicPr>
                <pic:cNvPr id="{doc_pr_id}" name="{name}"/>
                <pic:cNvPicPr/>
              </pic:nvPicPr>
              <pic:blipFill>
                <a:blip r:embed="{rId}"/>
                <a:stretch><a:fillRect/></a:stretch>
              </pic:blipFill>
              <pic:spPr>
                <a:xfrm>
                  <a:off x="0" y="0"/>
                  <a:ext cx="{cx}" cy="{cy}"/>
                </a:xfrm>
                <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
              </pic:spPr>
            </pic:pic>
          </a:graphicData>
        </a:graphic>
      </wp:inline>
    </w:drawing>
  </w:r>
</w:p>'''


def build_docx(source_docx, output_docx, image_map):
    """Replace placeholder tags in the source .docx with actual images."""
    print(f"=== Building .docx ===")

    with tempfile.TemporaryDirectory() as tmpdir:
        tmpdir = Path(tmpdir)

        # Extract the source docx
        with zipfile.ZipFile(source_docx, "r") as zf:
            zf.extractall(tmpdir)

        # Create media directory
        media_dir = tmpdir / "word" / "media"
        media_dir.mkdir(exist_ok=True)

        # Read document.xml as text (to find and replace placeholders)
        doc_xml_path = tmpdir / "word" / "document.xml"
        with open(doc_xml_path, "r", encoding="utf-8") as f:
            doc_xml = f.read()

        # Read rels file
        rels_path = tmpdir / "word" / "_rels" / "document.xml.rels"
        with open(rels_path, "r", encoding="utf-8") as f:
            rels_xml = f.read()

        # Read content types
        ct_path = tmpdir / "[Content_Types].xml"
        with open(ct_path, "r", encoding="utf-8") as f:
            ct_xml = f.read()

        # Add PNG content type if not present
        if 'Extension="png"' not in ct_xml:
            ct_xml = ct_xml.replace(
                "</Types>",
                '  <Default Extension="png" ContentType="image/png"/>\n</Types>'
            )
            with open(ct_path, "w", encoding="utf-8") as f:
                f.write(ct_xml)

        rId_base = 20
        replaced = 0

        for tag, img_path in sorted(image_map.items()):
            num = int(tag.split("_")[1])
            rId = f"rId{rId_base + num - 1}"
            media_name = img_path.name

            # Copy image to media folder
            shutil.copy2(img_path, media_dir / media_name)

            # Add relationship
            rel_entry = f'<Relationship Id="{rId}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="media/{media_name}"/>'
            rels_xml = rels_xml.replace("</Relationships>", f"  {rel_entry}\n</Relationships>")

            # Build the image paragraph XML
            drawing_xml = build_drawing_xml(
                rId=rId,
                doc_pr_id=20 + num,
                name=f"image_{num:02d}",
                desc=tag,
                cx=IMG_EMU_W,
                cy=IMG_EMU_H,
            )

            # Find and replace the placeholder paragraph
            # The placeholder looks like: <w:p>...<w:t>[IMAGE_XX]</w:t>...</w:p>
            pattern = rf'<w:p>.*?\[{tag}\].*?</w:p>'
            match = re.search(pattern, doc_xml, re.DOTALL)
            if match:
                doc_xml = doc_xml[:match.start()] + drawing_xml + doc_xml[match.end():]
                replaced += 1
                print(f"  Replaced [{tag}] with {media_name}")
            else:
                print(f"  Warning: Placeholder [{tag}] not found in document.xml")

        # Write modified files
        with open(doc_xml_path, "w", encoding="utf-8") as f:
            f.write(doc_xml)
        with open(rels_path, "w", encoding="utf-8") as f:
            f.write(rels_xml)

        # Repack as .docx
        output_docx.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(output_docx, "w", zipfile.ZIP_DEFLATED) as zf:
            for file_path in tmpdir.rglob("*"):
                if file_path.is_file():
                    arcname = file_path.relative_to(tmpdir)
                    zf.write(file_path, arcname)

        print(f"  Saved: {output_docx} ({replaced} images inserted)\n")
        return replaced


def build_pptx(prompts_file, images_dir, output_pptx):
    """Build a PowerPoint presentation with one image per slide."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    print(f"=== Building .pptx ===")

    with open(prompts_file, "r", encoding="utf-8") as f:
        prompts = json.load(f)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.333), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "UIAO"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "Unified Identity-Addressing-Overlay"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(0x5A, 0x5A, 0x5A)
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    p3.text = "Executive Brief  \u2022  April 2026"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0x5A, 0x5A, 0x5A)
    p3.alignment = PP_ALIGN.CENTER

    # Add cover image if available
    cover_img = images_dir / "image_01_cover_hero.png"
    if cover_img.exists():
        slide.shapes.add_picture(str(cover_img), Inches(3), Inches(4), Inches(7.333), Inches(3))

    # Section titles mapped to image IDs
    sections = [
        {"title": "The Problem", "image_id": "image_02", "body": "Every federal agency running M365 in GCC-Moderate faces the same grind. Manual compliance is slow, expensive, and stale before the ink dries."},
        {"title": "Continuous Monitoring", "image_id": "image_03", "body": "UIAO is designed to watch your environment continuously, check it against 323 FedRAMP Moderate Rev 5 controls, and detect drift within minutes."},
        {"title": "The Effort Comparison", "image_id": "image_04", "body": "The single largest cost in federal compliance is labor. UIAO is designed to dramatically reduce that category of work."},
        {"title": "Three Layers of Rules", "image_id": "image_05", "body": "FedRAMP Moderate Rev 5. CISA SCuBA & BOD 25-01. Your agency policies. Every configuration checked against all three."},
        {"title": "Immutable Evidence", "image_id": "image_06", "body": "Every check, every result, every timestamp is stored in an immutable chain with cryptographic signatures."},
        {"title": "What Makes This Different", "image_id": "image_07", "body": "Deterministic. Produces actual deliverables. Built specifically for federal M365."},
        {"title": "The Ask: 90-Day Pilot", "image_id": "image_08", "body": "Zero licensing cost. Open source. Connect, validate, champion. Real-world validation in a production tenant."},
    ]

    for section in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = section["title"]
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0x2E, 0x75, 0xB6)

        # Image
        img_file = None
        for prompt in prompts:
            if prompt["id"] == section["image_id"]:
                img_file = images_dir / prompt["filename"]
                break

        if img_file and img_file.exists():
            slide.shapes.add_picture(str(img_file), Inches(0.5), Inches(1.3), Inches(6), Inches(5.5))

        # Body text
        txBox = slide.shapes.add_textbox(Inches(7), Inches(1.3), Inches(5.833), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section["body"]
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x2D, 0x2D, 0x2D)

    prs.save(str(output_pptx))
    print(f"  Saved: {output_pptx} ({len(sections) + 1} slides)\n")


def build_pdf(input_docx, output_pdf):
    """Convert .docx to .pdf using LibreOffice."""
    print(f"=== Building .pdf ===")

    output_dir = output_pdf.parent

    # Try LibreOffice
    for soffice in ["soffice", "libreoffice", r"C:\Program Files\LibreOffice\program\soffice.exe",
                     r"C:\Program Files (x86)\LibreOffice\program\soffice.exe"]:
        try:
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(output_dir), str(input_docx)],
                capture_output=True, text=True, timeout=120
            )
            if result.returncode == 0:
                # LibreOffice names the output based on input filename
                generated = output_dir / (input_docx.stem + ".pdf")
                if generated.exists() and generated != output_pdf:
                    shutil.move(str(generated), str(output_pdf))
                print(f"  Saved: {output_pdf}\n")
                return True
        except (FileNotFoundError, subprocess.TimeoutExpired):
            continue

    print("  Warning: LibreOffice not found. Skipping PDF generation.")
    print("  Install LibreOffice to enable PDF output.\n")
    return False


def build_epub(prompts_file, images_dir, source_docx, output_epub):
    """Build an EPUB from the document content and images."""
    from ebooklib import epub

    print(f"=== Building .epub ===")

    # Extract text from source docx for content
    with open(prompts_file, "r", encoding="utf-8") as f:
        prompts = json.load(f)

    book = epub.EpubBook()
    book.set_identifier("uiao-executive-brief-2026")
    book.set_title("UIAO Executive Brief")
    book.set_language("en")
    book.add_author("UIAO Modernization Program")

    # Add images to the book
    epub_images = {}
    for item in prompts:
        img_path = images_dir / item["filename"]
        if img_path.exists():
            with open(img_path, "rb") as f:
                img_data = f.read()
            epub_img = epub.EpubImage()
            epub_img.file_name = f"images/{item['filename']}"
            epub_img.media_type = "image/png"
            epub_img.content = img_data
            book.add_item(epub_img)
            epub_images[item["id"]] = f"images/{item['filename']}"

    # Extract text content from the docx
    try:
        import zipfile
        from lxml import etree as ET

        with zipfile.ZipFile(source_docx, "r") as zf:
            doc_xml = zf.read("word/document.xml")

        tree = ET.fromstring(doc_xml)
        ns = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}

        paragraphs = []
        for p in tree.findall(".//w:p", ns):
            texts = []
            for t in p.findall(".//w:t", ns):
                if t.text:
                    texts.append(t.text)
            text = "".join(texts).strip()
            if text:
                # Check if it is a heading
                style = p.find(".//w:pStyle", ns)
                style_val = style.get(f"{{{ns['w']}}}val") if style is not None else None
                paragraphs.append({"text": text, "style": style_val})
    except Exception as e:
        print(f"  Warning: Could not extract text from docx: {e}")
        paragraphs = []

    # Build chapters by section
    sections = []
    current_section = {"title": "UIAO Executive Brief", "content": [], "image_id": "image_01"}

    section_image_map = {
        "The Problem": "image_02",
        "What UIAO Does About It": "image_03",
        "The Effort, Time, and Money Conversation": "image_04",
        "How It Works (Without the Jargon)": "image_05",
        "What Makes This Different": "image_07",
        "What We Need From Leadership": "image_08",
    }

    for para in paragraphs:
        if para["style"] and "Heading" in para["style"]:
            if current_section["content"] or current_section["title"]:
                sections.append(current_section)
            current_section = {
                "title": para["text"],
                "content": [],
                "image_id": section_image_map.get(para["text"], None),
            }
        elif para["text"].startswith("[IMAGE_"):
            continue  # Skip placeholders
        else:
            current_section["content"].append(para["text"])

    sections.append(current_section)

    # Build EPUB chapters
    chapters = []
    for i, section in enumerate(sections):
        chapter = epub.EpubHtml(
            title=section["title"],
            file_name=f"chapter_{i:02d}.xhtml",
            lang="en",
        )

        html = f'<h1 style="color: #2E75B6;">{section["title"]}</h1>\n'

        # Add image if mapped
        if section.get("image_id") and section["image_id"] in epub_images:
            html += f'<p style="text-align: center;"><img src="{epub_images[section["image_id"]]}" alt="{section["title"]}" style="max-width: 100%;"/></p>\n'

        for para_text in section["content"]:
            html += f"<p>{para_text}</p>\n"

        chapter.content = html
        book.add_item(chapter)
        chapters.append(chapter)

    # Table of contents and spine
    book.toc = [(epub.Section("Executive Brief"), chapters)]
    book.add_item(epub.EpubNcx())
    book.add_item(epub.EpubNav())
    book.spine = ["nav"] + chapters

    # Add basic CSS
    style = epub.EpubItem(
        uid="style",
        file_name="style/default.css",
        media_type="text/css",
        content=b"body { font-family: Georgia, serif; line-height: 1.6; margin: 1em; } h1 { color: #2E75B6; } img { max-width: 100%; height: auto; }",
    )
    book.add_item(style)

    epub.write_epub(str(output_epub), book, {})
    print(f"  Saved: {output_epub}\n")


def main():
    parser = argparse.ArgumentParser(description="UIAO Executive Brief — Document Build Pipeline")
    parser.add_argument("--skip-images", action="store_true", help="Skip image generation, use existing")
    parser.add_argument("--formats", nargs="+", default=["docx", "pptx", "pdf", "epub"],
                        choices=["docx", "pptx", "pdf", "epub"], help="Output formats to generate")
    parser.add_argument("--dry-run", action="store_true", help="Show plan without executing")
    parser.add_argument("--force-regen", action="store_true", help="Force regenerate all images")
    args = parser.parse_args()

    print("=" * 60)
    print("  UIAO Executive Brief — Document Build Pipeline")
    print("=" * 60)
    print(f"  Source:  {SOURCE_DOCX}")
    print(f"  Prompts: {PROMPTS_FILE}")
    print(f"  Images:  {IMAGES_DIR}")
    print(f"  Formats: {', '.join(args.formats)}")
    print("=" * 60)

    if args.dry_run:
        print("\n[DRY RUN] Would perform these steps:")
        if not args.skip_images:
            print("  1. Generate images from prompts.json")
        for fmt in args.formats:
            print(f"  - Build .{fmt} output")
        return

    # Step 1: Generate images
    if not args.skip_images:
        if args.force_regen:
            # Delete existing images to force regeneration
            for f in IMAGES_DIR.glob("*.png"):
                f.unlink()

        if not PROMPTS_FILE.exists():
            print(f"Error: Prompts file not found: {PROMPTS_FILE}")
            sys.exit(1)

        generate_images(PROMPTS_FILE, IMAGES_DIR, API_KEY, MODEL)

    # Build image map
    image_map = get_image_map(PROMPTS_FILE, IMAGES_DIR)
    if not image_map:
        print("Error: No images available. Run without --skip-images first.")
        sys.exit(1)

    # Step 2: Build outputs
    if "docx" in args.formats:
        build_docx(SOURCE_DOCX, OUTPUT_DOCX, image_map)

    if "pptx" in args.formats:
        try:
            build_pptx(PROMPTS_FILE, IMAGES_DIR, OUTPUT_PPTX)
        except ImportError:
            print("  Warning: python-pptx not installed. Run: pip install python-pptx\n")

    if "pdf" in args.formats:
        # PDF is generated from the .docx output
        if OUTPUT_DOCX.exists():
            build_pdf(OUTPUT_DOCX, OUTPUT_PDF)
        else:
            print("  Warning: .docx must be built first for PDF conversion\n")

    if "epub" in args.formats:
        try:
            build_epub(PROMPTS_FILE, IMAGES_DIR, SOURCE_DOCX, OUTPUT_EPUB)
        except ImportError:
            print("  Warning: ebooklib not installed. Run: pip install ebooklib\n")

    print("=" * 60)
    print("  Build complete!")
    print("=" * 60)
    for fmt in args.formats:
        output = {"docx": OUTPUT_DOCX, "pptx": OUTPUT_PPTX, "pdf": OUTPUT_PDF, "epub": OUTPUT_EPUB}[fmt]
        status = "OK" if output.exists() else "MISSING"
        print(f"  [{status}] {output}")
    print()


if __name__ == "__main__":
    main()
    main()
